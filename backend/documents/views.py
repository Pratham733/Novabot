from rest_framework import generics, permissions
from rest_framework.response import Response
from rest_framework.decorators import api_view, permission_classes
from services.ai import chat_complete
from .models import Document, ConvertedFile
from .serializers import DocumentSerializer
from typing import TYPE_CHECKING

if TYPE_CHECKING:
	# Help the type checker know common Document attributes without importing at runtime
	from .typings import Document as _DocumentType  # type: ignore
import os
import logging
import traceback

logger = logging.getLogger(__name__)


class DocumentListCreateView(generics.ListCreateAPIView):
	serializer_class = DocumentSerializer
	permission_classes = [permissions.IsAuthenticated]

	def get_queryset(self):
		return Document.objects.filter(owner=self.request.user)


class DocumentDetailView(generics.RetrieveUpdateDestroyAPIView):
	serializer_class = DocumentSerializer
	permission_classes = [permissions.IsAuthenticated]

	def get_queryset(self):
		return Document.objects.filter(owner=self.request.user)


class ConvertedFileListView(generics.ListAPIView):
	permission_classes = [permissions.IsAuthenticated]

	def get_queryset(self):
		return ConvertedFile.objects.filter(owner=self.request.user).order_by('-created_at')

	def list(self, request, *args, **kwargs):
		qs = self.get_queryset()
		data = [
			{
				'id': obj.id,
				'original_name': obj.original_name,
				'target_format': obj.target_format,
				'created_at': obj.created_at.isoformat(),
				'download_url': obj.file.url if obj.file else None,
			}
			for obj in qs
		]
		return Response({'results': data})


@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def regenerate_document(request, pk: int):
	try:
		doc = Document.objects.get(pk=pk, owner=request.user)
	except Document.DoesNotExist:
		return Response({'detail': 'Not found'}, status=404)

	instructions = request.data.get('instructions') or ''
	provider = request.data.get('provider') or 'openai'
	model = request.data.get('model')
	temperature = request.data.get('temperature') or 0.7

	system_prompt = (
		"You are an assistant that (re)generates high-quality {doc_type} content. "
		"Follow the user's instructions and keep the tone professional."
	).format(doc_type=doc.doc_type.replace('_', ' '))
	user_prompt = (
		f"Regenerate the following document titled '{doc.title}'. "
		f"Incorporate these instructions if provided: {instructions}\n\n"
		f"Current content:\n{doc.content}\n\nReturn only the regenerated content."
	)
	messages = [
		{"role": "system", "content": system_prompt},
		{"role": "user", "content": user_prompt},
	]

	ai_resp = chat_complete(messages, model=model, temperature=temperature, provider=provider)
	if ai_resp.get('error'):
		return Response({'error': ai_resp['error']}, status=502)

	new_content = ai_resp.get('content') or doc.content

	prev_versions = doc.meta.get('history', [])
	prev_versions.append({'content': doc.content, 'finalized': doc.meta.get('finalized', False)})
	doc.meta['history'] = prev_versions[-5:]
	doc.meta['finalized'] = False

	doc.content = new_content
	doc.save()
	return Response(DocumentSerializer(doc).data)


@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def finalize_document(request, pk: int):
	try:
		doc = Document.objects.get(pk=pk, owner=request.user)
	except Document.DoesNotExist:
		return Response({'detail': 'Not found'}, status=404)

	model = request.data.get('model')
	provider = request.data.get('provider') or 'openai'
	temperature = request.data.get('temperature') or 0.4
	messages = [
		{"role": "system", "content": "You are an assistant that polishes user documents for professionalism and clarity."},
		{"role": "user", "content": f"Polish this {doc.doc_type} titled '{doc.title}'. Return only the improved content.\n\n{doc.content}"}
	]
	ai_resp = chat_complete(messages, model=model, temperature=temperature, provider=provider)
	if ai_resp.get('error'):
		return Response({'error': ai_resp['error']}, status=502)

	new_content = ai_resp.get('content') or doc.content
	prev_versions = doc.meta.get('history', [])
	prev_versions.append({'content': doc.content, 'finalized': False})
	doc.meta['history'] = prev_versions[-5:]
	doc.content = new_content
	doc.meta['finalized'] = True
	doc.save()
	return Response(DocumentSerializer(doc).data)


@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def generate_document(request):
	doc_type = request.data.get('doc_type') or 'custom'
	# Normalize unsupported incoming types to closest allowed choice
	allowed = {c[0] for c in getattr(Document, 'TYPE_CHOICES', [])}
	fallback_map = {
		'proposal': 'proposal',
		'email': 'email',
		'summary': 'summary',
		'presentation': 'presentation',
		'contract': 'contract',
	}
	if doc_type not in allowed:
		# Try mapped fallback else default to 'custom'
		doc_type = fallback_map.get(doc_type, 'custom')
	title = request.data.get('title') or ''
	prompt = request.data.get('prompt') or ''
	provider = request.data.get('provider') or 'openai'
	model = request.data.get('model')
	temperature = request.data.get('temperature') or 0.7

	system_prompt = (
		"You are an assistant that generates high-quality {doc_type} content. "
		"Follow the user's prompt and keep the tone professional."
	).format(doc_type=doc_type.replace('_', ' '))

	user_prompt = f"{prompt}\n\nReturn only the generated content."
	messages = [
		{"role": "system", "content": system_prompt},
		{"role": "user", "content": user_prompt},
	]

	ai_resp = chat_complete(messages, model=model, temperature=temperature, provider=provider)
	if ai_resp.get('error'):
		return Response({'error': ai_resp['error']}, status=502)

	content = ai_resp.get('content') or ai_resp.get('text') or prompt

	doc = Document.objects.create(owner=request.user, doc_type=doc_type, title=title, content=content or '')
	return Response(DocumentSerializer(doc).data, status=201)


@api_view(['GET'])
@permission_classes([permissions.IsAuthenticated])
def export_document(request, pk: int):
	try:
		doc = Document.objects.get(pk=pk, owner=request.user)
	except Document.DoesNotExist:
		return Response({'detail': 'Not found'}, status=404)

	# Help the type checker know that doc has an id attribute for filename formatting
	doc_id = getattr(doc, 'id', None)

	fmt = (request.GET.get('format') or 'txt').lower()
	serializer = DocumentSerializer(doc)

	if fmt == 'json':
		resp = Response(serializer.data)
		resp['Content-Disposition'] = f'attachment; filename="document-{doc_id}.json"'
		return resp

	if fmt == 'txt':
		content = doc.content or ''
		return Response(content, content_type='text/plain; charset=utf-8', headers={'Content-Disposition': f'attachment; filename="document-{doc_id}.txt"'})

	if fmt == 'pdf':
		try:
			from io import BytesIO
			from reportlab.pdfgen import canvas
		except Exception:
			return Response({'error': 'PDF export not available: install reportlab'}, status=501)
		buffer = BytesIO()
		p = canvas.Canvas(buffer)
		y = 800
		for line in (doc.content or '').splitlines() or ['']:
			if y < 40:
				p.showPage()
				y = 800
			p.drawString(40, y, line)
			y -= 14
		p.save()
		buffer.seek(0)
		return Response(buffer.getvalue(), content_type='application/pdf', headers={'Content-Disposition': f'attachment; filename="document-{doc_id}.pdf"'})

	if fmt == 'docx':
		try:
			from io import BytesIO
			import docx
		except Exception:
			return Response({'error': 'DOCX export not available: install python-docx'}, status=501)
		buffer = BytesIO()
		d = docx.Document()
		for line in (doc.content or '').splitlines():
			d.add_paragraph(line)
		d.save(buffer)
		buffer.seek(0)
		return Response(buffer.getvalue(), content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document', headers={'Content-Disposition': f'attachment; filename="document-{doc_id}.docx"'})

	if fmt == 'pptx':
		try:
			from io import BytesIO
			from pptx import Presentation
		except Exception:
			return Response({'error': 'PPTX export not available: install python-pptx'}, status=501)
		prs = Presentation()
		slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
		for line in (doc.content or '').split('\n\n'):
			slide = prs.slides.add_slide(slide_layout)
			title = slide.shapes.title
			title.text = line[:50]
		buffer = BytesIO()
		prs.save(buffer)
		buffer.seek(0)
		return Response(buffer.getvalue(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', headers={'Content-Disposition': f'attachment; filename="document-{doc_id}.pptx"'})

	return Response({'error': 'Format not supported'}, status=400)


@api_view(['POST'])
@permission_classes([permissions.IsAuthenticated])
def convert_document(request):
	"""Convert an uploaded file to a desired format. Supports txt/json and attempts pdf/docx/pptx when libs are installed."""
	upload = request.FILES.get('file')
	target = (request.POST.get('format') or 'txt').lower()
	persist = request.POST.get('persist') == 'true'
	if not upload:
		return Response({'error': 'No file uploaded'}, status=400)
	if target not in {'txt','json','pdf','docx','pptx'}:
		return Response({'error': f'Unsupported target format: {target}'}, status=400)

	logger.info('convert_document called: filename=%s, target=%s, size=%s', getattr(upload, 'name', None), target, getattr(upload, 'size', None))

	try:
		content_bytes = upload.read()
	except Exception as e:
		logger.exception('Failed to read uploaded file')
		return Response({'error': 'Failed to read uploaded file', 'detail': str(e)}, status=500)

	# If DOCX, try to extract text first
	filename = getattr(upload, 'name', '') or 'upload'
	_, ext = os.path.splitext(filename.lower())
	if ext == '.docx':
		try:
			from io import BytesIO
			import docx as _docx
			docx_file = _docx.Document(BytesIO(content_bytes))
			paragraphs = [p.text for p in docx_file.paragraphs]
			extracted = '\n\n'.join([p for p in paragraphs if p])
			if extracted:
				content_bytes = extracted.encode('utf-8')
		except Exception:
			logger.exception('DOCX text extraction failed; using raw bytes')

	# Helper to optionally persist bytes
	def _persist_bytes(data: bytes, out_ext: str, original_name: str):
		if not persist or not request.user.is_authenticated:
			return None
		from django.core.files.base import ContentFile
		cf = ContentFile(data)
		try:
			stored = ConvertedFile.objects.create(
				owner=request.user,
				original_name=original_name,
				target_format=out_ext.lstrip('.'),
				file=cf,
				meta={'retention_hours': 24}
			)
			stored.file.save(f"converted-{stored.pk}.{out_ext.lstrip('.')}", cf, save=True)
			return stored
		except Exception:
			return None

	# txt
	if target == 'txt':
		try:
			text = content_bytes.decode('utf-8')
		except Exception:
			text = content_bytes.decode('latin1', errors='ignore')
		data = text.encode('utf-8')
		stored = _persist_bytes(data, 'txt', filename)
		resp = Response(text, content_type='text/plain; charset=utf-8', headers={'Content-Disposition': f'attachment; filename="converted-{filename}.txt"'})
		if stored:
			resp.data = {'id': stored.id, 'download_url': stored.file.url, 'filename': stored.file.name}
		return resp

	# json
	if target == 'json':
		try:
			text = content_bytes.decode('utf-8')
		except Exception:
			text = content_bytes.decode('latin1', errors='ignore')
		data = text.encode('utf-8')
		stored = _persist_bytes(data, 'json', filename)
		payload = {'filename': filename, 'content': text}
		if stored:
			payload.update({'id': stored.id, 'download_url': stored.file.url})
		return Response(payload)

	# pdf
	if target == 'pdf':
		try:
			from io import BytesIO
			from reportlab.pdfgen import canvas
		except Exception as e:
			logger.exception('PDF conversion library missing')
			return Response({'error': 'PDF conversion not available: install reportlab', 'detail': str(e)}, status=501)
		buffer = BytesIO()
		p = canvas.Canvas(buffer)
		y = 800
		try:
			text = content_bytes.decode('utf-8')
		except Exception:
			text = content_bytes.decode('latin1', errors='ignore')
		for line in text.splitlines() or ['']:
			if y < 40:
				p.showPage()
				y = 800
			p.drawString(40, y, line)
			p.drawString(40, y, line)
			y -= 14
		p.save()
		buffer.seek(0)
		pdf_bytes = buffer.getvalue()
		stored = _persist_bytes(pdf_bytes, 'pdf', filename)
		resp = Response(pdf_bytes, content_type='application/pdf', headers={'Content-Disposition': f'attachment; filename="converted-{filename}.pdf"'})
		if stored:
			resp['X-Converted-Id'] = str(stored.id)
		return resp

	# docx
	if target == 'docx':
		try:
			from io import BytesIO
			import docx
		except Exception as e:
			logger.exception('DOCX conversion library missing')
			return Response({'error': 'DOCX conversion not available: install python-docx', 'detail': str(e)}, status=501)
		buffer = BytesIO()
		try:
			text = content_bytes.decode('utf-8')
		except Exception:
			text = content_bytes.decode('latin1', errors='ignore')
		doc = docx.Document()
		for line in text.splitlines():
			doc.add_paragraph(line)
		doc.save(buffer)
		buffer.seek(0)
		docx_bytes = buffer.getvalue()
		stored = _persist_bytes(docx_bytes, 'docx', filename)
		resp = Response(docx_bytes, content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document', headers={'Content-Disposition': f'attachment; filename="converted-{filename}.docx"'})
		if stored:
			resp['X-Converted-Id'] = str(stored.id)
		return resp

	# pptx
	if target == 'pptx':
		try:
			from io import BytesIO
			from pptx import Presentation
		except Exception as e:
			logger.exception('PPTX conversion library missing')
			return Response({'error': 'PPTX conversion not available: install python-pptx', 'detail': str(e)}, status=501)
		prs = Presentation()
		slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
		try:
			text = content_bytes.decode('utf-8')
		except Exception:
			text = content_bytes.decode('latin1', errors='ignore')
		for block in text.split('\n\n'):
			slide = prs.slides.add_slide(slide_layout)
			if slide.shapes.title:
				slide.shapes.title.text = (block or '')[:200]
		buffer = BytesIO()
		prs.save(buffer)
		buffer.seek(0)
		pptx_bytes = buffer.getvalue()
		stored = _persist_bytes(pptx_bytes, 'pptx', filename)
		resp = Response(pptx_bytes, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', headers={'Content-Disposition': f'attachment; filename="converted-{filename}.pptx"'})
		if stored:
			resp['X-Converted-Id'] = str(stored.id)
		return resp

	# unsupported
	return Response({'error': 'Conversion for requested format not available on server. Install appropriate libraries (python-docx, reportlab, python-pptx) to enable.'}, status=501)




@api_view(['GET'])
@permission_classes([permissions.AllowAny])
def convert_capabilities(request):
	"""Report which conversion formats are available on the server."""
	formats = ['txt', 'json']
	# optional libraries
	try:
		import reportlab  # type: ignore
		formats.append('pdf')
	except Exception:
		pass
	try:
		import docx as _docx  # type: ignore
		formats.append('docx')
	except Exception:
		pass
	try:
		import pptx as _pptx  # type: ignore
		formats.append('pptx')
	except Exception:
		pass
	return Response({'formats': formats})
